#!/usr/bin/env python
# -*- coding: utf-8 -*-
# @Date  : 2021/4/28 5:44 下午
# @File  : translate_ppt.py
# @Author: johnson
# @Desc  : m2m100的翻译的API接口

from flask import Flask, request, jsonify, abort
import os
from pptx import Presentation
from transformers import M2M100ForConditionalGeneration, M2M100Tokenizer
import logging
import re
import torch
import gc
import atexit,time
from apscheduler.schedulers.background import BackgroundScheduler
import pandas as pd
import pymysql

def unload_model_schedule():
    global nobusy_count
    global model
    if model.model:
        nobusy_count += 1
        app.logger.info(f"Model 是存在的，不使用GPU的时间是 {nobusy_count*10}秒")
    if nobusy_count >= 30:
        #开始清理GPU, 30次，一共300秒，即5分钟不使用GPU，就卸载
        app.logger.info(f"开始清理模型")
        model.unload_model()
        nobusy_count = 0

scheduler = BackgroundScheduler()
scheduler.add_job(func=unload_model_schedule, trigger="interval", seconds=10)
scheduler.start()

atexit.register(lambda: scheduler.shutdown())

app = Flask(__name__)
app.config['DEBUG'] = False
if os.path.exists('/data/var/log/'):
    app.config['API_LOG_FILE'] = '/data/var/log/translate_api.log'
else:
    app.config['API_LOG_FILE'] = 'translate_api.log'
# 日志配置信息, Running on之类的显示在日志里面
if app.config['DEBUG']:
    logging.basicConfig(filename=app.config['API_LOG_FILE'], level=logging.DEBUG)
else:
    logging.basicConfig(filename=app.config['API_LOG_FILE'], level=logging.INFO)

class TranslateModel(object):
    def __init__(self, verbose=False):
        self.model_name = "./translate_model"
        self.excel_file = '翻译对照表总表.xlsx'
        # 判断使用的设备
        self.device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
        self.model = None
        self.tokenizer = None
        # 提取中英文的替换的关键字列表
        self.corres_keys = None
        self.cn_keys = None
        self.en_keys = None
        self.connet_mysql()
        # beam size 为3
        self.num_beams = 3
        # 返回句子为3
        self.num_return_sequences = 3
        # 防止重复生成单词的ngram
        self.no_repeat_ngram_size = 2
    def load_keywords_from_mysql(self):
        """
        不同于extract_dict从excel中加载，这里从mysql数据库加载
        :return:
        :rtype:
        """
        self.connet_mysql()
        sql = 'select * from pro_translate_dict'
        df = pd.read_sql(sql,self.db)
        self.cn_keys = df['cnword'].tolist()
        self.en_keys = df['enword'].tolist()
        self.corres_keys = df['modelword'].tolist()

    def connet_mysql(self):
        db = pymysql.connect(host="192.168.xx.xx",
                             user="xx",
                             password="xxx",
                             port=3306,
                             db='xxx',
                             charset='utf8mb4',
                             )
        self.db = db
    def submit_mysql(self,cn_keys,en_keys,corres_keys):
        # 提交数据data到数据库表pro_translate_dict
        total = 0
        loss = 0
        for cnword, enword, modelword in zip(cn_keys,en_keys,corres_keys):
            cursor = self.db.cursor()
            sql = "INSERT INTO `pro_translate_dict` (`cnword`, `enword`, `modelword`) VALUES ('%s','%s', '%s')" % (
            cnword, enword, modelword)
            try:
                cursor.execute(sql)
                self.db.commit()
                total += 1
                app.logger.info(f"提交{cnword},{enword},{modelword}数据到mysql成功")
            except Exception as e:
                self.db.rollback()
                app.logger.info(f"提交{cnword},{enword},{modelword}数据到mysql失败")
                print(e)
                loss += 1
        self.db.commit()
        app.logger.info(f"提交{total}条数据到mysql成功, {loss}条失败")

    def load_model(self):
        """
        加载模型
        :return:
        """
        app.logger.info(f"开始加载模型")
        model = M2M100ForConditionalGeneration.from_pretrained(self.model_name)
        model.to(self.device)
        self.tokenizer = M2M100Tokenizer.from_pretrained(self.model_name)
        self.model = model

    def extract_dict(self, excel_file, write2mysql=False):
        """
        提取excel文件中的中英文关键字
        :param excel_file: excel文件
        :param write2mysql: 是否把结果写到mysql
        :return:
        :rtype:
        """
        df = pd.read_excel(excel_file)
        cn_keywords = df['中文'].tolist()
        en_keywords = df['英文'].tolist()
        app.logger.info(f"处理前中文和英文单词个数分别是 {len(cn_keywords)}, {len(en_keywords)}")
        cn_keys = []
        en_keys = []
        # 处理一下/分隔的词，这样的词是多个的
        for cn, en in zip(cn_keywords, en_keywords):
            if '/' in cn:
                cn_words = cn.split('/')
            else:
                cn_words = [cn]
            if '/' in en:
                en_words = en.split('/')
            else:
                en_words = [en]
            # 中英文单词都加入单词表
            for c in cn_words:
                for e in en_words:
                    cn_keys.append(c)
                    en_keys.append(e)
        app.logger.info(f"处理后中文和英文单词个数分别是 {len(cn_keys)}, {len(en_keys)}")
        app.logger.info(f"中文关键字cn_keys是{cn_keys}")
        app.logger.info(f"英文关键字en_keys是{en_keys}")
        self.cn_keys = cn_keys
        self.en_keys = en_keys
        # 每个中文关键字对应的翻译结果，从模型中获取
        app.logger.info(f"首先处理对应的关键字的中文到英文的翻译结果")
        corres_keys = []
        for cnkey in cn_keys:
            #每个中文对应的英文翻译结果
            corres_key = self.translate2en(text=cnkey,do_replace=False)
            corres_keys.append(corres_key)
        app.logger.info(f"corres_keys是 {corres_keys}")
        self.corres_keys = corres_keys
        if write2mysql:
            #开始同步到mysql数据库
            self.submit_mysql(cn_keys,en_keys,corres_keys)

    def start_replace(self, cntext, entext):
        """
        使用cn_keys和en_keys进行正则替换
        :param cntext:
        :param entext:
        :return:
        """
        result = entext
        for cnkey, enkey, corres_key in zip(self.cn_keys, self.en_keys, self.corres_keys):
            if cnkey in cntext:
                result = re.sub(corres_key,enkey,entext,flags=re.I)
                if result != entext:
                    app.logger.info(f"进行了替换: {entext},被从{corres_key}替换成{enkey}")
        return result
    def translate2en(self, text, do_replace=True):
        """
        翻译中文text到英文
        :param do_replace: 默认使用excel中的关键字进行替换
        """
        global nobusy_count
        nobusy_count = 0
        if self.model is None:
            #如果predict_model没有加载，自动加载默认的模型
            self.load_model()
        res = re.findall('[\u4e00-\u9fa5]+', text)
        if not res:
            app.logger.info(f"原文是: {text}, 不包含中文字符，不需要翻译")
            return text
        self.tokenizer.src_lang = "zh"
        encoded_zh = self.tokenizer(text, return_tensors="pt")
        encoded_zh.data['attention_mask'] = encoded_zh.data['attention_mask'].to(self.device)
        encoded_zh.data['input_ids'] = encoded_zh.data['input_ids'].to(self.device)
        generated_tokens = self.model.generate(**encoded_zh, num_beams=self.num_beams, num_return_sequences=self.num_return_sequences,no_repeat_ngram_size=self.no_repeat_ngram_size, forced_bos_token_id=self.tokenizer.get_lang_id("en"))
        entext = self.tokenizer.batch_decode(generated_tokens, skip_special_tokens=True)
        app.logger.info(f"原文是: {text}\n模型翻译的结果是: {entext}")
        result_text = entext[0]
        for ent in entext:
            #如果没发现中文，就用这个翻译结果，如果发现中文，就换一个翻译结果，如果仍然有中文，那也没办法了
            if not re.findall('[\u4e00-\u9fa5]+', ent):
                result_text = ent
                break
        if do_replace:
            result_text = self.start_replace(cntext=text, entext=result_text)
        app.logger.info(f"原文是: {text}\n最终翻译的结果是: {entext}")
        return result_text
    def unload_model(self):
        """
        卸载模型，不把模型加载到GPU中
        :return:
        :rtype:
        """
        app.logger.info(f"开始卸载模型")
        self.model = None
        self.tokenizer = None
        gc.collect()
        torch.cuda.empty_cache()
        app.logger.info(f"模型卸载完成")

def slade_change(SRC_PPT, TARGET_PPT, DICT_FILE=None):
    """
    中文到英文PPT的翻译，SRC_PPT中文PPT，TARGET_PPT英文PPT
    :param DICT_FILE: excel格式的，包含中文和英文2列，用于正则替换
    """
    if DICT_FILE:
        # 如果传了excel，那么用用户传的excel做替换，否则用默认的
        model.extract_dict(excel_file=DICT_FILE)
    else:
        model.load_keywords_from_mysql()
    prs = Presentation(SRC_PPT)
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            # 翻译表格
            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        cell_text = cell.text_frame.text
                        cell_entext = model.translate2en(cell_text)
                        prs.slides[slide_idx].shapes[shape_idx].table.rows[row_idx].cells[cell_idx].text_frame.text = cell_entext
            # 图表翻译
            if shape.has_chart:
                # 图表的标题翻译
                title_txt = shape.chart.chart_title.text_frame.text
                title_entext = model.translate2en(title_txt)
                shape.chart.chart_title.text_frame.text = title_entext
            #翻译其它文本
            if not shape.has_text_frame:
                continue
            for paragraph_idx,paragraph in enumerate(shape.text_frame.paragraphs):
                paragraph_text = ""
                for run_idx, run in enumerate(paragraph.runs):
                    paragraph_text = paragraph_text + run.text
                    # 只要第一个run的txt，其它的都设为空
                    if run_idx != 0:
                        prs.slides[slide_idx].shapes[shape_idx].text_frame.paragraphs[paragraph_idx].runs[run_idx].text = ''
                if paragraph_text:
                    entext = model.translate2en(paragraph_text)
                    prs.slides[slide_idx].shapes[shape_idx].text_frame.paragraphs[paragraph_idx].runs[0].text = entext
    prs.save(TARGET_PPT)
    app.logger.info(f"读取{SRC_PPT},修改完成了{len(prs.slides)}页的PPT{TARGET_PPT}")
    return 200, "成功"

@app.route("/fileTranslate", methods=['POST'])
def translate():
    """
    翻译api
    :return:
    """
    form_dict = request.form.to_dict()
    json_dict = request.get_json()
    app.logger.info(f"用户请求的form内容是{form_dict}")
    app.logger.info(f"用户请求的json内容是{json_dict}")
    # dict_file是翻译后检查和替换的字典
    if json_dict:
        source_ppt = json_dict.get('inputFilePath', None)
        des_ppt = json_dict.get('resultFilePath', None)
        dict_file = json_dict.get('dictFilePath', None)
    else:
        source_ppt = form_dict.get('inputFilePath', None)
        des_ppt = form_dict.get('resultFilePath', None)
        dict_file = form_dict.get('dictFilePath', None)
    if source_ppt.split('.')[-1].lower() not in ['pptx']:
        app.logger.warning('ppt格式不符合要求')
        return jsonify({'ret': -101, 'msg': 'ppt格式不符合要求'})
    if not os.path.exists(source_ppt):
        app.logger.warning('ppt文件不存在，请检查服务器上是否存在这个ppt')
        return jsonify({'ret': -102, 'msg': 'ppt文件不存在，请检查服务器上是否存在这个ppt'})
    if dict_file and not os.path.exists(dict_file):
        app.logger.warning('给了dictFilePath参数，但文件不在服务器上')
        return jsonify({'ret': -103, 'msg': '给了dictFilePath参数，但文件不在服务器上'})
    # 保存文件
    code, msg = slade_change(SRC_PPT=source_ppt, TARGET_PPT=des_ppt, DICT_FILE=dict_file)
    # 判断OCR后端识别的结果是否正确
    if code == 200:
        return jsonify({'ret': 0, 'msg': '成功'})
    else:
        return jsonify({'ret': -104, 'msg': msg})


@app.route("/syncMysql", methods=['POST'])
def syncmysql():
    """
    把dictfile的内容写到mysql
    :return:
    """
    form_dict = request.form.to_dict()
    json_dict = request.get_json()
    app.logger.info(f"准备同步excel内容到mysql数据库")
    app.logger.info(f"用户请求的form内容是{form_dict}")
    app.logger.info(f"用户请求的json内容是{json_dict}")
    # dict_file是翻译后检查和替换的字典
    if json_dict:
        source_ppt = json_dict.get('inputFilePath', None)
        des_ppt = json_dict.get('resultFilePath', None)
        dict_file = json_dict.get('dictFilePath', None)
    else:
        source_ppt = form_dict.get('inputFilePath', None)
        des_ppt = form_dict.get('resultFilePath', None)
        dict_file = form_dict.get('dictFilePath', None)
    if dict_file and not os.path.exists(dict_file):
        app.logger.warning('给了dictFilePath参数，但文件不在服务器上')
        return jsonify({'ret': -103, 'msg': '给了dictFilePath参数，但文件不在服务器上'})
    # 保存文件
    model.extract_dict(excel_file=dict_file,write2mysql=True)
    return jsonify({'ret': 0, 'msg': '成功'})

if __name__ == "__main__":
    # 预训练模型
    model = TranslateModel()
    nobusy_count = 0
    # slade_change(SRC_PPT="图表翻译.pptx",TARGET_PPT="表格翻译-英文.pptx")
    app.run(host='0.0.0.0', port=3325, debug=False, threaded=True)
